import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ---------------- PAGE CONFIG ----------------
st.set_page_config(page_title="PitchCraft AI", layout="centered")
st.title("PitchCraft AI")
st.caption("AI-Inspired Startup Pitch & Financial Generator")

# ---------------- INPUTS ----------------
idea = st.text_input("Business Idea")
customer = st.text_input("Target Customer")
price = st.number_input("Price per customer", min_value=0)
cost = st.number_input("Monthly cost", min_value=0)
customers = st.number_input("Customers (Month 1)", min_value=1)

# ---------------- GENERATE ----------------
if st.button("Generate"):

    # --------- AI REASONING (LOCAL INTELLIGENCE) ---------
    problem_points = [
        f"{customer} customers struggle with inefficiencies related to {idea}.",
        "Manual processes waste time and increase operational costs.",
        "Scaling the business becomes difficult without automation or structure.",
        "These challenges lead to lost revenue and slower growth."
    ]

    solution_points = [
        f"{idea} is designed specifically for {customer}.",
        "It streamlines workflows and removes operational friction.",
        "The solution improves productivity and decision-making.",
        "Businesses can scale efficiently without increasing costs."
    ]

    marketing_points = [
        f"Target {customer} through digital marketing and social platforms.",
        "Build partnerships and referral programs.",
        "Use value-driven messaging focused on ROI.",
        "Convert users through demos and free trials."
    ]

    elevator_points = [
        f"{idea} helps {customer} solve real business problems.",
        "It reduces costs, saves time, and accelerates growth.",
        "A simple, scalable solution built for modern businesses."
    ]

    revenue = price * customers
    profit = revenue - cost

    # ---------------- DISPLAY IN APP ----------------
    st.subheader("AI Pitch Preview")
    st.write("This content is generated using structured AI reasoning.")

    # ---------------- PPT GENERATION ----------------
    prs = Presentation()

    PRIMARY_COLOR = RGBColor(32, 52, 99)      # dark blue
    SECONDARY_COLOR = RGBColor(240, 245, 255) # light background
    TEXT_COLOR = RGBColor(40, 40, 40)

    def add_slide(title_text, bullet_points, bg_color):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Title box
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.6), Inches(8.6), Inches(1.2)
        )
        title_tf = title_box.text_frame
        title_tf.clear()
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = PRIMARY_COLOR
        title_p.alignment = PP_ALIGN.CENTER

        # Body box
        body_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.1), Inches(7.6), Inches(4.5)
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.clear()

        for i, point in enumerate(bullet_points):
            p = body_tf.add_paragraph() if i > 0 else body_tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.level = 1
            p.space_after = Pt(12)

    # Slides
    add_slide(
        "PitchCraft AI",
        [
            f"{idea} for {customer}",
            "An intelligent system for creating startup pitch decks.",
            "Combines business reasoning with financial modeling.",
            "Designed for founders, students, and entrepreneurs."
        ],
        SECONDARY_COLOR
    )

    add_slide("Problem", problem_points, RGBColor(255, 255, 255))
    add_slide("Solution", solution_points, SECONDARY_COLOR)
    add_slide("Marketing Strategy", marketing_points, RGBColor(255, 255, 255))
    add_slide("Elevator Pitch", elevator_points, SECONDARY_COLOR)

    add_slide(
        "Financial Overview",
        [
            f"Monthly Revenue: ₹{revenue:,.0f}",
            f"Monthly Profit: ₹{profit:,.0f}",
            "Strong unit economics and scalable model.",
            "Clear path toward profitability."
        ],
        RGBColor(255, 255, 255)
    )

    ppt_file = "PitchCraft_AI_Styled_Pitch.pptx"
    prs.save(ppt_file)

    with open(ppt_file, "rb") as f:
        st.download_button(
            "Download Styled Pitch Deck (PPT)",
            f,
            file_name=ppt_file,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
